import os
from docx import Document
from docx.shared import Pt, Cm
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.text import MSO_VERTICAL_ANCHOR
import locale
locale.setlocale(locale.LC_ALL, 'de_DE.UTF-8')

s3ll = os.listdir("./qrcodes/s3ll")
s3ll_list = [x for x in s3ll if "QR" in x]
s3ll_list = sorted(s3ll_list, key=locale.strxfrm)

pll = os.listdir("./qrcodes/pll")
pll_list = [x for x in pll if "QR" in x]
pll_list.sort()

create_pll = False
create_s3ll = True

rotation_left = 5400000
rotation_right = 16200000
xpos_left = Cm(2.9)
xpos_right = Cm(12.85)
ypos_start = Cm(0)
shape_type = 152
shape_width = Cm(3.3)
shape_height = Cm(9.1)
y_spacer = Cm(0.5)
y_spacer_text = Cm(0.49)
tf_width = Cm(6.35)
tf_height = Cm(2.70)
tf_start_y = Cm(3.2)
tf_spacer_y = Cm(1.1)
qrcode_y_start = Cm(3.02)
qrcode_size = Cm(3.09)
qrcode_spacer_y = Cm(0.7)
page_max_y = Cm(28.5)

y_spacer_buffer = Cm(0.5)
y_factor = (y_spacer_buffer + y_spacer) / y_spacer
y_spacer += y_spacer_buffer

slide_idx = 0

def calculate_top_position(is_left_shape, xpos_left, xpos_right, new_y_pos, shape_width, shape_height):
    # Check the shape position and set initial parameters
    if is_left_shape:
        top = new_y_pos - shape_width  # Since the shape is rotated, the top position will shift by the width of the shape
    else:
        top = new_y_pos + shape_height  # Since the shape is rotated, the top position will shift by the height of the shape

    return top

if create_s3ll:
    prs = Presentation("./qrcodes/s3ll/S3LL Übersicht_Template.pptx")
    left_shape_template = None
    right_shape_template = None
    for shape in prs.slides[0].shapes:
        try:
            xfrm = shape._element.spPr.get_or_add_xfrm()
            print(shape.auto_shape_type)
            print(int(xfrm.get("rot")))
            pass
        except:
            pass

    target_slide = prs.slides[slide_idx]
    new_y_pos = ypos_start
    is_left_shape = False

    counter = 0
    for s3ll in s3ll_list:
        is_left_shape = not is_left_shape

        new_shape_bottom = new_y_pos + shape_height
        if is_left_shape and new_shape_bottom >= page_max_y:
            slide_idx += 1
            target_slide = prs.slides[slide_idx]
            new_y_pos = ypos_start
            is_left_shape = True
            print("New page index: %s" % slide_idx)
            counter = 0

        if is_left_shape:
            left = xpos_left
            frame_rotation = str(270*60000)
            rotation = 90
            text_y_pos = new_y_pos - shape_width
        else:
            left = xpos_right
            frame_rotation = str(90*60000)
            rotation = 270
            text_y_pos = new_y_pos + shape_height

        new_shape = target_slide.shapes.add_shape(
            shape_type,
            left,
            new_y_pos,
            shape_width,
            shape_height
        )

        new_shape.fill.solid()
        new_shape.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
        new_shape.rotation = rotation

        new_shape.line.color.rgb = new_shape.fill.fore_color.rgb

        if is_left_shape:
            text_x_pos = 0
            qr_code_x = Cm(5.84)
        else:
            text_x_pos = Cm(9.95)
            qr_code_x = text_x_pos + Cm(5.84)

        text_y_pos = counter * Cm(4.3) + Cm(2.9)
        tf_height = shape_width

        text_shape = target_slide.shapes.add_shape(
            1,
            text_x_pos,
            text_y_pos,
            tf_width,
            tf_height
        )
        text_shape.fill.background()
        text_shape.line.fill.background()

        para = text_shape.text_frame.paragraphs[-1]
        r = para.add_run()
        r.text = s3ll.replace("QR_", "").split(".")[0]
        r.font.color.rgb = RGBColor(0, 0, 0)
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.name = "Lucida Sans"

        para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        text_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        qr_code_y = text_y_pos + Pt(4)

        qr_code = target_slide.shapes.add_picture("./qrcodes/s3ll/%s" % s3ll, qr_code_x, qr_code_y, width=qrcode_size, height=qrcode_size)

        if not is_left_shape:
            new_y_pos += y_spacer + shape_width
            counter += 1

    prs.save("Test.pptx")

if create_pll:
    img_width = 80
    row = 0
    column = 0
    maxcolumns = 5

    doc = Document("./qrcodes/pll/PLL Uebersicht_Template.docx")
    table = doc.tables[0]

    for pll in pll_list:
        cell = table.cell(row, column)
        paragraph = cell.paragraphs[0]
        run = paragraph.add_run()
        img_path = "./qrcodes/pll/%s" % pll.replace("QR_", "").replace(".png", ".jpg")
        img = run.add_picture(img_path)
        aspect = img.height/img.width
        img.width = Pt(img_width)
        img.height = int(aspect * img.width)

        cell = table.cell(row + 1, column)
        paragraph = cell.paragraphs[-1]
        run = paragraph.add_run()
        run.text = pll.replace("QR_", "").replace(".png", "")

        cell = table.cell(row + 2, column)
        paragraph = cell.paragraphs[0]
        run = paragraph.add_run()
        img_path = "./qrcodes/pll/%s" % pll
        img = run.add_picture(img_path)
        aspect = img.height / img.width
        img.width = Pt(img_width)
        img.height = int(aspect * img.width)

        column += 1
        if column >= maxcolumns:
            row += 4
            column = 0

    doc.save("./Test.docx")